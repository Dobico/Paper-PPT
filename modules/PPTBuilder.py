from pptx import Presentation
import os

def create_summary_from_template(summary_text: str,
                                  template_path: str = None,
                                  style_option: str = "default",
                                  output_path: str = "summary_output.pptx"):
    # 1. 프레젠테이션 불러오기
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    layout_map = {
        "default": 1,
        "title_only": 5,
        "section_header": 2,
        "blank": 6
    }
    slide_layout_index = layout_map.get(style_option, 1)

    # 2. 요약 내용 분할
    sections = summary_text.strip().split("\n")
    slides = []

    current_title = None
    current_content = ""

    for line in sections:
        line = line.strip()
        if line.startswith(("1.", "2.", "3.", "4.")):
            if current_title and current_content:
                slides.append((current_title, current_content.strip()))
            parts = line.split(" ", 1)
            current_title = parts[1] if len(parts) > 1 else line
            current_content = ""
        else:
            current_content += " " + line

    if current_title and current_content:
        slides.append((current_title, current_content.strip()))

    # 3. 슬라이드 생성
    for title, content in slides:
        try:
            slide_layout = prs.slide_layouts[slide_layout_index]
        except IndexError:
            slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(slide_layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 1:
                placeholder.text = content
                break

    prs.save(output_path)
    print(f"✅ 슬라이드 저장 완료: {output_path}")